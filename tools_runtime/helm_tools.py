#!/usr/bin/env python3
"""External open-source tool adapter used by Helm 1.9.1.

Each subcommand writes a compact UTF-8 result to stdout and diagnostics to
stderr. The C++ host owns timeouts/cancellation and never executes arbitrary
shell strings.
"""
from __future__ import annotations

import argparse
import json
import os
import random
import shutil
import subprocess
import sys
import time
import urllib.parse
from pathlib import Path


def fail(message: str, code: int = 2) -> None:
    print(message, file=sys.stderr)
    raise SystemExit(code)


def _configure_playwright() -> None:
    if os.environ.get("PLAYWRIGHT_BROWSERS_PATH"):
        return
    try:
        root = Path(sys.executable).resolve().parents[2]
    except IndexError:
        return
    candidate = root / "Playwright-Browsers"
    if candidate.exists():
        os.environ["PLAYWRIGHT_BROWSERS_PATH"] = str(candidate)


def _compact_text(text: str, max_chars: int) -> str:
    return " ".join((text or "").replace("\x00", " ").split())[:max_chars]


def _validate_web_url(url: str) -> str:
    parsed = urllib.parse.urlparse(url.strip())
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        raise ValueError("URL must be an absolute http or https address")
    return parsed.geturl()


def _extract_html_text(html: str, max_chars: int) -> str:
    text = ""
    try:
        import trafilatura
        text = trafilatura.extract(
            html,
            include_comments=False,
            include_tables=True,
            include_links=False,
            favor_recall=True,
        ) or ""
    except Exception:
        text = ""
    if len(text.strip()) < 200:
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html, "html.parser")
            for node in soup(["script", "style", "noscript", "svg", "template"]):
                node.decompose()
            text = soup.get_text(" ", strip=True)
        except Exception:
            pass
    return _compact_text(text, max_chars)


class _BrowserPool:
    """One Chromium launch shared by every fallback in a single tool call.

    The previous shape launched a whole browser per URL. Three fallbacks in one
    search meant three launches, three teardowns, and three full timeout budgets
    stacked end to end — which is what made a search look like it had hung.
    """

    def __init__(self) -> None:
        self._pw = None
        self._browser = None
        self.error = ""

    def available(self) -> bool:
        return self.error == ""

    def _ensure(self):
        if self._browser is not None:
            return self._browser
        _configure_playwright()
        try:
            from playwright.sync_api import sync_playwright
        except Exception as exc:
            self.error = f"Playwright is not installed: {exc}"
            return None
        try:
            self._pw = sync_playwright().start()
            self._browser = self._pw.chromium.launch(headless=True)
        except Exception as exc:
            # Almost always the Chromium binaries were never downloaded.
            self.error = (f"headless Chromium could not start ({exc}). "
                          "Run install_helm_tools.cmd to download it.")
            self.close()
            return None
        return self._browser

    def fetch(self, url: str, max_chars: int, budget_seconds: float = 25.0) -> dict:
        browser = self._ensure()
        if browser is None:
            return {"url": url, "text": "", "method": "browser", "error": self.error}
        page = None
        try:
            page = browser.new_page(
                viewport={"width": 1440, "height": 1000},
                user_agent=("Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                            "HelmLocalAI/1.9.1.0 Chrome/126 Safari/537.36"),
            )
            # Budget split: most of it on first paint, a short tail for late XHR.
            page.goto(url, wait_until="domcontentloaded", timeout=int(budget_seconds * 1000 * 0.7))
            try:
                page.wait_for_load_state("networkidle", timeout=int(budget_seconds * 1000 * 0.2))
            except Exception:
                pass
            page.wait_for_timeout(250)
            title = page.title()
            final_url = page.url
            try:
                body_text = page.locator("body").inner_text(timeout=4_000)
            except Exception:
                body_text = _extract_html_text(page.content(), max_chars)
            return {
                "url": final_url or url,
                "title": _compact_text(title, 500),
                "text": _compact_text(body_text, max_chars),
                "method": "browser",
                "error": "",
            }
        except Exception as exc:
            return {"url": url, "text": "", "method": "browser", "error": str(exc)}
        finally:
            if page is not None:
                try:
                    page.close()
                except Exception:
                    pass

    def close(self) -> None:
        try:
            if self._browser is not None:
                self._browser.close()
        except Exception:
            pass
        try:
            if self._pw is not None:
                self._pw.stop()
        except Exception:
            pass
        self._browser = None
        self._pw = None


def _needs_browser(text: str) -> bool:
    """True when a static fetch clearly did not get the real content."""
    lowered = text.lower()
    markers = ("enable javascript", "javascript is required", "checking your browser",
               "please turn on javascript", "requires javascript")
    return len(text) < 500 or any(marker in lowered for marker in markers)


def _static_fetch(url: str, max_chars: int, client=None, timeout: float = 12.0) -> dict:
    try:
        url = _validate_web_url(url)
    except ValueError as exc:
        return {"url": url, "text": "", "method": "static", "error": str(exc)}
    try:
        import httpx
    except Exception as exc:
        return {"url": url, "text": "", "method": "static", "error": f"httpx unavailable: {exc}"}
    headers = {
        "User-Agent": ("Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                       "HelmLocalAI/1.9.1.0 Chrome/126 Safari/537.36"),
        "Accept": "text/html,application/xhtml+xml,text/plain;q=0.9,*/*;q=0.6",
        "Accept-Language": "en-US,en;q=0.8",
    }
    own = client is None
    http = client or httpx.Client(headers=headers, follow_redirects=True, timeout=timeout)
    try:
        response = http.get(url, timeout=timeout)
        response.raise_for_status()
        content_type = response.headers.get("content-type", "").lower()
        text = ""
        if any(kind in content_type for kind in ("html", "xml", "text")) or not content_type:
            text = _extract_html_text(response.text, max_chars)
        title = ""
        if response.text and ("html" in content_type or "<html" in response.text[:1000].lower()):
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(response.text, "html.parser")
                title = _compact_text(soup.title.get_text(" ", strip=True), 500) if soup.title else ""
            except Exception:
                title = ""
        return {"url": str(response.url), "title": title, "text": text,
                "method": "static", "error": ""}
    except Exception as exc:
        return {"url": url, "text": "", "method": "static", "error": str(exc)}
    finally:
        if own:
            http.close()


def fetch_page(url: str, max_chars: int = 12000, client=None, pool=None,
               browser_budget: float = 25.0) -> dict:
    """Static fetch, escalating to a real browser only when the page needs one."""
    result = _static_fetch(url, max_chars, client)
    if result.get("text") and not _needs_browser(result["text"]):
        return result

    own_pool = pool is None
    pool = pool or _BrowserPool()
    try:
        rendered = pool.fetch(url, max_chars, browser_budget)
    finally:
        if own_pool:
            pool.close()

    if rendered.get("text"):
        return rendered
    # Both paths failed: keep whatever static managed, and report both reasons so
    # a missing Chromium is visible instead of looking like an empty page.
    combined = "; ".join(part for part in (result.get("error", ""), rendered.get("error", "")) if part)
    result["error"] = combined or "no readable text found"
    result["browser_attempted"] = True
    return result


def web_fetch(args: argparse.Namespace) -> None:
    print(json.dumps(fetch_page(args.url, max(1000, min(args.max_chars, 100000))),
                     ensure_ascii=False))


# How many search results get opened, and how many of those may escalate to a
# browser. Both are deliberately small: every browser fallback is seconds the
# user spends watching a progress bar.
_AUTO_OPEN = 3
_MAX_BROWSER_FALLBACKS = 2
_SEARCH_DEADLINE_SECONDS = 75.0


def seen_list(args: argparse.Namespace) -> None:
    """Disk-backed dedup for perpetual agent runs.

    An agent that clears context between batches has no memory of what it already
    processed, so it would resurvey the same repositories forever. This keeps the
    set of processed identifiers in a file the agent checks and appends to. Two
    modes: --add writes identifiers, default reads them back.
    """
    path = Path(args.file)
    path.parent.mkdir(parents=True, exist_ok=True)
    existing = set()
    if path.exists():
        existing = {l.strip() for l in path.read_text(encoding="utf-8", errors="replace").splitlines() if l.strip()}

    if args.add:
        added = [x.strip() for x in args.add.split(",") if x.strip() and x.strip() not in existing]
        if added:
            with path.open("a", encoding="utf-8") as f:
                for x in added:
                    f.write(x + "\n")
        print(json.dumps({"added": added, "total_seen": len(existing) + len(added)}, ensure_ascii=False))
        return

    print(json.dumps({"total_seen": len(existing),
                      "seen": sorted(existing)[:args.limit],
                      "note": "Repositories already processed in earlier batches. Skip any repo whose "
                              "full name appears here; do not survey it again."}, ensure_ascii=False))


def github_search(args: argparse.Namespace) -> None:
    """Search GitHub repositories through the REST API.

    Returns fully-formed records: real url, stars, license, pushed_at, language,
    archived flag, open issue count, and the description. This exists because a
    model told to "find repos and record their details" will otherwise invent
    plausible-looking github.com/<owner>/<name> URLs from training memory and
    fetch 404s. Here there is nothing to invent - the API is the source.
    """
    try:
        import httpx
    except Exception as exc:
        fail(f"github search dependency missing: {exc}")

    per_page = max(1, min(args.max_results, 50))
    sort = args.sort if args.sort in ("stars", "updated", "forks", "help-wanted-issues") else "best-match"
    params = {"q": args.query, "per_page": per_page}
    if sort != "best-match":
        params["sort"] = sort
    params["order"] = "desc"

    headers = {
        "Accept": "application/vnd.github+json",
        "X-GitHub-Api-Version": "2022-11-28",
        "User-Agent": "HelmLocalAI/1.9.1",
    }
    # An optional token lifts the rate limit from 10/min to 30/min. Read-only.
    token = os.environ.get("GITHUB_TOKEN", "").strip()
    if token:
        headers["Authorization"] = f"Bearer {token}"

    try:
        with httpx.Client(headers=headers, follow_redirects=True, timeout=25) as client:
            resp = client.get("https://api.github.com/search/repositories", params=params)
            if resp.status_code == 403 and "rate limit" in resp.text.lower():
                fail("github rate limit reached (unauthenticated is 10 searches/min). "
                     "Wait a minute, or set GITHUB_TOKEN. This is not a failure of the query.")
            resp.raise_for_status()
            payload = resp.json()
    except SystemExit:
        raise
    except Exception as exc:
        fail(f"github search failed: {exc}")

    repos = []
    for item in payload.get("items", []):
        lic = item.get("license") or {}
        repos.append({
            "name": item.get("full_name", ""),
            "url": item.get("html_url", ""),
            "description": _compact_text(item.get("description") or "", 300),
            "stars": item.get("stargazers_count", 0),
            "language": item.get("language") or "unknown",
            "license": lic.get("spdx_id") or lic.get("name") or "none",
            "pushed_at": (item.get("pushed_at") or "")[:10],
            "open_issues": item.get("open_issues_count", 0),
            "archived": bool(item.get("archived", False)),
            "is_fork": bool(item.get("fork", False)),
            "topics": item.get("topics", [])[:8],
        })

    print(json.dumps({
        "query": args.query,
        "total_found": payload.get("total_count", 0),
        "returned": len(repos),
        "repositories": repos,
        "note": ("These are verified GitHub API results. Use the url and other fields exactly as "
                 "given - do NOT construct or guess repository URLs. Every field here is authoritative; "
                 "only fetch a page if you need README detail the description does not cover."),
    }, ensure_ascii=False))


def web_search(args: argparse.Namespace) -> None:
    try:
        import httpx
        from bs4 import BeautifulSoup
    except Exception as exc:
        fail(f"web search dependencies missing: {exc}")

    started = time.time()
    headers = {
        "User-Agent": ("Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                       "HelmLocalAI/1.9.1.0 Chrome/126 Safari/537.36"),
        "Accept-Language": "en-US,en;q=0.8",
    }
    url = "https://html.duckduckgo.com/html/?" + urllib.parse.urlencode({"q": args.query})
    requested = max(1, min(args.max_results, 12))
    results: list[dict] = []
    pool = _BrowserPool()

    try:
        with httpx.Client(headers=headers, follow_redirects=True, timeout=15) as client:
            response = client.get(url)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, "html.parser")
            for row in soup.select(".result"):
                link = row.select_one("a.result__a")
                if not link:
                    continue
                href = link.get("href", "")
                parsed = urllib.parse.urlparse(href)
                if parsed.netloc.endswith("duckduckgo.com"):
                    href = urllib.parse.parse_qs(parsed.query).get("uddg", [href])[0]
                if not href.startswith(("http://", "https://")):
                    continue
                snippet = row.select_one(".result__snippet")
                results.append({
                    "title": " ".join(link.get_text(" ", strip=True).split()),
                    "url": href,
                    "snippet": " ".join(snippet.get_text(" ", strip=True).split()) if snippet else "",
                })
                if len(results) >= max(requested, _AUTO_OPEN):
                    break

            targets = results[:_AUTO_OPEN]

            # Static fetches run concurrently. They are almost all network wait,
            # so three at once costs about what one used to.
            from concurrent.futures import ThreadPoolExecutor
            if targets:
                with ThreadPoolExecutor(max_workers=len(targets)) as pool_exec:
                    statics = list(pool_exec.map(
                        lambda r: _static_fetch(str(r["url"]), 9000, timeout=12.0), targets))
            else:
                statics = []

            for result, fetched in zip(targets, statics):
                result["page_text"] = fetched.get("text", "")
                result["fetch_method"] = fetched.get("method", "")
                if fetched.get("error"):
                    result["fetch_error"] = fetched["error"]
                if fetched.get("title") and not result.get("title"):
                    result["title"] = fetched["title"]
                if fetched.get("url"):
                    result["final_url"] = fetched["url"]

            # Only pages that actually need JavaScript escalate, capped, and only
            # while there is time left in the budget.
            escalations = 0
            for result in targets:
                if escalations >= _MAX_BROWSER_FALLBACKS:
                    break
                remaining = _SEARCH_DEADLINE_SECONDS - (time.time() - started)
                if remaining < 12:
                    result.setdefault("fetch_note", "skipped browser rendering: time budget spent")
                    continue
                if not _needs_browser(str(result.get("page_text", ""))):
                    continue
                escalations += 1
                rendered = pool.fetch(str(result["url"]), 9000, min(25.0, remaining - 4))
                if rendered.get("text"):
                    result["page_text"] = rendered["text"]
                    result["fetch_method"] = "browser"
                    result.pop("fetch_error", None)
                    if rendered.get("url"):
                        result["final_url"] = rendered["url"]
                elif rendered.get("error"):
                    result["fetch_error"] = rendered["error"]
    finally:
        pool.close()

    usable = sum(1 for r in results[:_AUTO_OPEN] if len(str(r.get("page_text", ""))) >= 300)
    payload = {
        "query": args.query,
        "requested_max_results": requested,
        "results": results[:requested] if requested >= len(results[:_AUTO_OPEN]) else results,
        "auto_fetched_pages": usable,
        "elapsed_seconds": round(time.time() - started, 1),
        "javascript_rendering": "available" if pool.available() else "unavailable",
        "research_note": (
            "Helm already opened the top results with static extraction and escalated to headless "
            "Chromium where the page needed JavaScript. Use page_text directly. If evidence is still "
            "incomplete, call fetch_web_page or search_web again with a refined query in the same "
            "turn; do not ask the user to approve read-only research."
        ),
    }
    if not pool.available() and pool.error:
        payload["javascript_rendering_error"] = pool.error
    print(json.dumps(payload, ensure_ascii=False))


def extract_document(args: argparse.Namespace) -> None:
    path = Path(args.path)
    if not path.is_file():
        fail("document does not exist")
    ext = path.suffix.lower()
    chunks: list[str] = []
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            for index, page in enumerate(PdfReader(str(path)).pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    chunks.append(f"\n--- Page {index} ---\n{text}")
        elif ext == ".docx":
            from docx import Document
            doc = Document(str(path))
            chunks.extend(p.text for p in doc.paragraphs if p.text.strip())
            for table in doc.tables:
                for row in table.rows:
                    chunks.append("\t".join(cell.text for cell in row.cells))
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True, data_only=True)
            for ws in wb.worksheets:
                chunks.append(f"\n--- Sheet: {ws.title} ---")
                for row in ws.iter_rows(values_only=True):
                    values = ["" if value is None else str(value) for value in row]
                    if any(values):
                        chunks.append("\t".join(values))
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            for index, slide in enumerate(prs.slides, 1):
                chunks.append(f"\n--- Slide {index} ---")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        chunks.append(text)
        else:
            fail(f"unsupported document type: {ext}")
    except SystemExit:
        raise
    except Exception as exc:
        fail(f"document extraction failed: {exc}")
    text = "\n".join(chunks)
    sys.stdout.write(text[: args.max_chars])


def desktop_screenshot(args: argparse.Namespace) -> None:
    try:
        import mss
        from PIL import Image
    except Exception as exc:
        fail(f"screenshot dependencies missing: {exc}")
    target = Path(args.output)
    target.parent.mkdir(parents=True, exist_ok=True)
    with mss.mss() as capture:
        monitor = capture.monitors[args.monitor if 0 <= args.monitor < len(capture.monitors) else 0]
        shot = capture.grab(monitor)
        Image.frombytes("RGB", shot.size, shot.rgb).save(target)
    print(str(target))


def desktop_click(args: argparse.Namespace) -> None:
    try:
        import pyautogui
    except Exception as exc:
        fail(f"desktop automation dependency missing: {exc}")
    pyautogui.click(args.x, args.y, clicks=args.clicks, interval=max(0, args.interval_ms) / 1000, button=args.button)
    print(f"clicked {args.button} at ({args.x}, {args.y}) x{args.clicks}")


def desktop_type(args: argparse.Namespace) -> None:
    try:
        import pyautogui
    except Exception as exc:
        fail(f"desktop automation dependency missing: {exc}")
    pyautogui.write(args.text, interval=max(0, args.interval_ms) / 1000)
    print(f"typed {len(args.text)} characters")


def desktop_hotkey(args: argparse.Namespace) -> None:
    try:
        import pyautogui
    except Exception as exc:
        fail(f"desktop automation dependency missing: {exc}")
    keys = [key.strip().lower() for key in args.keys.split(",") if key.strip()]
    if not keys:
        fail("no keys supplied")
    pyautogui.hotkey(*keys)
    print("pressed " + "+".join(keys))


def _replace_workflow_values(workflow: dict, prompt: str, negative: str, width: int, height: int, prefix: str) -> None:
    text_nodes = []
    for node in workflow.values():
        if not isinstance(node, dict):
            continue
        class_type = str(node.get("class_type", ""))
        inputs = node.setdefault("inputs", {})
        if class_type == "CLIPTextEncode" and "text" in inputs:
            text_nodes.append(node)
        elif class_type in {"EmptyLatentImage", "EmptySD3LatentImage"}:
            inputs["width"] = width
            inputs["height"] = height
        elif class_type == "KSampler":
            inputs["seed"] = random.SystemRandom().randint(0, 2**63 - 1)
        elif class_type == "SaveImage":
            inputs["filename_prefix"] = prefix
    if text_nodes:
        text_nodes[0]["inputs"]["text"] = prompt
    if len(text_nodes) > 1:
        text_nodes[1]["inputs"]["text"] = negative


def _ensure_comfyui(base: str, start_command: str, wait_seconds: int = 180) -> None:
    import httpx

    def ready() -> bool:
        try:
            response = httpx.get(base + "/system_stats", timeout=3)
            return response.status_code == 200
        except Exception:
            return False

    if ready():
        return
    start_path = Path(start_command) if start_command else None
    if not start_path or not start_path.is_file():
        fail("ComfyUI is not running and START_COMFYUI.cmd was not found; run Install / repair tools")

    # Launch ComfyUI's own interpreter directly rather than going through
    # cmd.exe on the .cmd file. Routing through cmd gives the server a console
    # window, and START_COMFYUI.cmd ends in `pause`, so that window sticks
    # around after the server stops. CREATE_NO_WINDOW keeps it headless.
    root = start_path.parent
    venv_python = root / "ComfyUI" / ".venv" / "Scripts" / "python.exe"
    comfy_dir = root / "ComfyUI"
    if venv_python.is_file() and (comfy_dir / "main.py").is_file():
        command = [str(venv_python), "main.py", "--listen", "127.0.0.1", "--port", "8188"]
        workdir = str(comfy_dir)
    else:
        command = ["cmd.exe", "/d", "/c", str(start_path)]
        workdir = str(start_path.parent)

    try:
        flags = 0
        if os.name == "nt":
            flags = (getattr(subprocess, "CREATE_NO_WINDOW", 0)
                     | getattr(subprocess, "CREATE_NEW_PROCESS_GROUP", 0))
        subprocess.Popen(
            command,
            cwd=workdir,
            stdin=subprocess.DEVNULL,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            creationflags=flags,
        )
    except Exception as exc:
        fail(f"could not start ComfyUI: {exc}")

    deadline = time.time() + max(10, wait_seconds)
    while time.time() < deadline:
        if ready():
            return
        time.sleep(2)
    fail("ComfyUI did not become ready after it was started")


def comfy_generate(args: argparse.Namespace) -> None:
    try:
        import httpx
    except Exception as exc:
        fail(f"ComfyUI dependency missing: {exc}")
    workflow_path = Path(args.workflow)
    if not workflow_path.is_file():
        fail("ComfyUI API workflow JSON is not configured")
    try:
        workflow = json.loads(workflow_path.read_text(encoding="utf-8"))
    except Exception as exc:
        fail(f"invalid ComfyUI workflow JSON: {exc}")
    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    prefix = f"helm_{int(time.time())}"
    _replace_workflow_values(workflow, args.prompt, args.negative_prompt, args.width, args.height, prefix)

    base = args.url.rstrip("/")
    _ensure_comfyui(base, args.start_command)
    with httpx.Client(timeout=30) as client:
        response = client.post(base + "/prompt", json={"prompt": workflow, "client_id": "helm-local"})
        response.raise_for_status()
        prompt_id = response.json().get("prompt_id")
        if not prompt_id:
            fail("ComfyUI did not return a prompt id")
        deadline = time.time() + args.timeout_seconds
        history = None
        while time.time() < deadline:
            history_response = client.get(base + f"/history/{prompt_id}")
            if history_response.status_code == 200:
                history = history_response.json().get(prompt_id)
                if history and history.get("outputs"):
                    break
            time.sleep(1)
        if not history or not history.get("outputs"):
            fail("ComfyUI generation timed out")

        saved: list[str] = []
        for output in history["outputs"].values():
            for image in output.get("images", []):
                params = {"filename": image["filename"], "subfolder": image.get("subfolder", ""), "type": image.get("type", "output")}
                data = client.get(base + "/view", params=params)
                data.raise_for_status()
                target = output_dir / Path(image["filename"]).name
                target.write_bytes(data.content)
                saved.append(str(target))
        if not saved:
            fail("ComfyUI completed without an image output")
        print(json.dumps({"prompt_id": prompt_id, "files": saved}, ensure_ascii=False))


# ---------------------------------------------------------------------------
# OSINT pack. Every source here is a public registry or transparency log.
# Keyless where possible; the two that are not (HIBP, OpenCorporates) read an
# env var and say so plainly rather than failing with an opaque 401.
# ---------------------------------------------------------------------------

OSINT_UA = "HelmLocalAI/1.9.7"


def _osint_client(headers: dict | None = None, timeout: int = 30):
    try:
        import httpx
    except Exception as exc:
        fail(f"osint dependency missing: {exc}")
    h = {"User-Agent": OSINT_UA, "Accept": "application/json"}
    if headers:
        h.update(headers)
    return httpx.Client(headers=h, follow_redirects=True, timeout=timeout)


def _sec_headers() -> dict:
    # SEC fair-access policy: a request with no User-Agent, or a generic library
    # default, gets 403 and can earn a ~10 minute IP block. Policy asks for an
    # identifying string plus a contact address.
    ua = os.environ.get("SEC_USER_AGENT", "").strip()
    if not ua:
        ua = f"{OSINT_UA} (set SEC_USER_AGENT to 'YourName your@email' per SEC fair-access policy)"
    return {"User-Agent": ua, "Accept-Encoding": "gzip, deflate"}


def _rdap_summary(payload: dict) -> dict:
    out = {}
    out["handle"] = payload.get("handle") or payload.get("ldhName") or ""
    if payload.get("ldhName"):
        out["name"] = payload["ldhName"]
    for key in ("startAddress", "endAddress", "ipVersion", "type", "country", "name"):
        if payload.get(key):
            out[key] = payload[key]
    out["status"] = payload.get("status", [])
    events = {}
    for ev in payload.get("events", []):
        action = ev.get("eventAction")
        if action:
            events[action] = ev.get("eventDate", "")
    if events:
        out["events"] = events
    ns = [n.get("ldhName", "") for n in payload.get("nameservers", []) if n.get("ldhName")]
    if ns:
        out["nameservers"] = ns
    entities = []
    for ent in payload.get("entities", []):
        roles = ent.get("roles", [])
        label = ""
        vcard = ent.get("vcardArray")
        if isinstance(vcard, list) and len(vcard) > 1:
            for field in vcard[1]:
                if isinstance(field, list) and len(field) > 3 and field[0] == "fn":
                    label = str(field[3])
                    break
        entities.append({"roles": roles, "name": label or ent.get("handle", "")})
    if entities:
        out["entities"] = entities
    if payload.get("secureDNS"):
        out["dnssec"] = payload["secureDNS"].get("delegationSigned", None)
    return out


def _is_ip(value: str) -> bool:
    import ipaddress
    try:
        ipaddress.ip_address(value.strip())
        return True
    except ValueError:
        return False


def rdap_lookup(args: argparse.Namespace) -> None:
    """Registration data for a domain or IP via the RDAP bootstrap.

    RDAP is the IETF/ICANN replacement for port-43 WHOIS: structured JSON, one
    endpoint, no per-registry parser. rdap.org redirects to the authoritative
    registry or RIR. Registrant contact fields are usually redacted under GDPR;
    that is the registry's choice, not a failure here.
    """
    target = args.target.strip().lower().rstrip(".")
    if not target:
        fail("rdap: empty target")
    kind = "ip" if _is_ip(target) else "domain"
    url = f"https://rdap.org/{kind}/{urllib.parse.quote(target)}"
    try:
        with _osint_client() as client:
            resp = client.get(url)
            if resp.status_code == 404:
                fail(f"rdap: no record for {target}. Some ccTLDs (.io, .dev among them) "
                     f"publish no RDAP server yet; that is a coverage gap, not a missing domain.")
            resp.raise_for_status()
            payload = resp.json()
    except SystemExit:
        raise
    except Exception as exc:
        fail(f"rdap lookup failed: {exc}")
    summary = _rdap_summary(payload)
    summary["query"] = target
    summary["query_type"] = kind
    print(json.dumps(summary, ensure_ascii=False, indent=1))


def _doh(client, name: str, rtype: str) -> list:
    resp = client.get("https://dns.google/resolve", params={"name": name, "type": rtype})
    resp.raise_for_status()
    body = resp.json()
    return [a.get("data", "") for a in body.get("Answer", []) if a.get("data")]


def dns_lookup(args: argparse.Namespace) -> None:
    """DNS records over DoH (dns.google). Keyless, no resolver config needed."""
    name = args.domain.strip().lower().rstrip(".")
    if not name:
        fail("dns: empty domain")
    types = [t.strip().upper() for t in args.types.split(",") if t.strip()]
    if not types:
        types = ["A"]
    results = {}
    try:
        with _osint_client() as client:
            for rtype in types[:8]:
                try:
                    results[rtype] = _doh(client, name, rtype)
                except Exception as exc:
                    results[rtype] = [f"error: {exc}"]
    except Exception as exc:
        fail(f"dns lookup failed: {exc}")
    print(json.dumps({"domain": name, "records": results}, ensure_ascii=False, indent=1))


def crtsh_subdomains(args: argparse.Namespace) -> None:
    """Subdomains from Certificate Transparency logs via crt.sh.

    CT logs record every issued certificate, so this surfaces hosts that were
    never published in DNS. crt.sh is free but frequently overloaded and has no
    documented rate limit; slow responses are normal, which is why the caller
    runs this as a job.
    """
    domain = args.domain.strip().lower().rstrip(".")
    if not domain:
        fail("crtsh: empty domain")
    url = "https://crt.sh/"
    params = {"q": f"%.{domain}", "output": "json"}
    try:
        with _osint_client(timeout=args.timeout) as client:
            resp = client.get(url, params=params)
            resp.raise_for_status()
            if not resp.text.strip():
                print(json.dumps({"domain": domain, "subdomains": [], "note": "crt.sh returned empty"}))
                return
            rows = resp.json()
    except Exception as exc:
        fail(f"crt.sh query failed: {exc}. The service is often overloaded; retry is usually enough.")
    found = set()
    for row in rows:
        for entry in str(row.get("name_value", "")).split("\n"):
            entry = entry.strip().lower().rstrip(".")
            if not entry or entry.startswith("*"):
                continue
            if entry == domain or entry.endswith("." + domain):
                found.add(entry)
    subs = sorted(found)
    limit = max(1, args.max_results)
    print(json.dumps({
        "domain": domain,
        "total_found": len(subs),
        "returned": min(len(subs), limit),
        "subdomains": subs[:limit],
    }, ensure_ascii=False, indent=1))


def _cik_for(client, query: str) -> tuple[str, str]:
    resp = client.get("https://www.sec.gov/files/company_tickers.json", headers=_sec_headers())
    resp.raise_for_status()
    table = resp.json()
    q = query.strip().lower()
    best = None
    for row in table.values():
        ticker = str(row.get("ticker", "")).lower()
        title = str(row.get("title", "")).lower()
        if ticker == q:
            best = row
            break
        if best is None and q in title:
            best = row
    if not best:
        fail(f"edgar: no filer matched '{query}'. Try the exact ticker.")
    return str(best["cik_str"]).zfill(10), best.get("title", "")


def edgar_company(args: argparse.Namespace) -> None:
    """Filer profile and recent filings from EDGAR. Keyless; UA header required."""
    try:
        with _osint_client(timeout=30) as client:
            cik, title = _cik_for(client, args.query)
            resp = client.get(f"https://data.sec.gov/submissions/CIK{cik}.json", headers=_sec_headers())
            if resp.status_code == 403:
                fail("edgar: 403. Set SEC_USER_AGENT to 'YourName your@email' - the SEC "
                     "rejects requests without an identifying User-Agent.")
            resp.raise_for_status()
            data = resp.json()
    except SystemExit:
        raise
    except Exception as exc:
        fail(f"edgar company lookup failed: {exc}")
    recent = data.get("filings", {}).get("recent", {})
    n = max(1, min(args.max_results, 40))
    filings = []
    for i in range(min(n, len(recent.get("form", [])))):
        filings.append({
            "form": recent["form"][i],
            "filed": recent["filingDate"][i],
            "accession": recent["accessionNumber"][i],
            "doc": (f"https://www.sec.gov/Archives/edgar/data/{int(cik)}/"
                    f"{recent['accessionNumber'][i].replace('-', '')}/"
                    f"{recent['primaryDocument'][i]}") if recent.get("primaryDocument") else "",
        })
    print(json.dumps({
        "cik": cik,
        "name": data.get("name", title),
        "tickers": data.get("tickers", []),
        "exchanges": data.get("exchanges", []),
        "sic": data.get("sicDescription", ""),
        "state": data.get("stateOfIncorporation", ""),
        "recent_filings": filings,
    }, ensure_ascii=False, indent=1))


def edgar_search(args: argparse.Namespace) -> None:
    """Full-text search across EDGAR filings since 2001 (efts.sec.gov)."""
    params = {"q": args.query, "from": 0}
    if args.forms.strip():
        params["forms"] = args.forms.strip()
    try:
        with _osint_client(timeout=40) as client:
            resp = client.get("https://efts.sec.gov/LATEST/search-index",
                              params=params, headers=_sec_headers())
            resp.raise_for_status()
            body = resp.json()
    except Exception as exc:
        fail(f"edgar full-text search failed: {exc}")
    hits = body.get("hits", {})
    out = []
    for hit in hits.get("hits", [])[:max(1, min(args.max_results, 30))]:
        src = hit.get("_source", {})
        adsh = src.get("adsh", "")
        out.append({
            "adsh": adsh,
            "form": src.get("file_type") or src.get("root_form", ""),
            "filed": src.get("file_date", ""),
            "entity": (src.get("display_names") or [""])[0],
            "id": hit.get("_id", ""),
        })
    print(json.dumps({
        "query": args.query,
        "total_hits": hits.get("total", {}).get("value", 0),
        "results": out,
    }, ensure_ascii=False, indent=1))


def hibp_account(args: argparse.Namespace) -> None:
    """Breaches for an account. Requires a paid HIBP subscription key."""
    key = os.environ.get("HIBP_API_KEY", "").strip()
    if not key:
        fail("hibp: no HIBP_API_KEY set. Every HIBP endpoint that searches by email or "
             "domain needs a paid subscription key (haveibeenpwned.com/API/Key). Only the "
             "Pwned Passwords range API is free, and it does not take an email address.")
    account = args.account.strip()
    if not account:
        fail("hibp: empty account")
    url = f"https://haveibeenpwned.com/api/v3/breachedaccount/{urllib.parse.quote(account)}"
    try:
        with _osint_client(headers={"hibp-api-key": key}, timeout=25) as client:
            resp = client.get(url, params={"truncateResponse": "false"})
            if resp.status_code == 404:
                print(json.dumps({"account": account, "breached": False, "breaches": []}))
                return
            if resp.status_code == 401:
                fail("hibp: 401 - key rejected. Keys are 32-char hex and tied to a live subscription.")
            if resp.status_code == 429:
                fail("hibp: 429 - rate limited by your subscription tier. Back off and retry.")
            resp.raise_for_status()
            rows = resp.json()
    except SystemExit:
        raise
    except Exception as exc:
        fail(f"hibp lookup failed: {exc}")
    breaches = [{
        "name": b.get("Name", ""),
        "domain": b.get("Domain", ""),
        "breach_date": b.get("BreachDate", ""),
        "accounts": b.get("PwnCount", 0),
        "data": b.get("DataClasses", []),
        "verified": b.get("IsVerified", None),
    } for b in rows]
    print(json.dumps({"account": account, "breached": True, "count": len(breaches),
                      "breaches": breaches}, ensure_ascii=False, indent=1))


def opencorporates_search(args: argparse.Namespace) -> None:
    """Company registry search. Requires an OpenCorporates API token."""
    token = os.environ.get("OPENCORPORATES_API_TOKEN", "").strip()
    if not token:
        fail("opencorporates: no OPENCORPORATES_API_TOKEN set. The API requires a token; "
             "it is free only for open-data projects that republish under a share-alike "
             "licence, otherwise it is a paid plan. See opencorporates.com.")
    params = {"q": args.query, "api_token": token, "per_page": max(1, min(args.max_results, 30))}
    if args.jurisdiction.strip():
        params["jurisdiction_code"] = args.jurisdiction.strip().lower()
    try:
        with _osint_client(timeout=30) as client:
            resp = client.get("https://api.opencorporates.com/v0.4/companies/search", params=params)
            if resp.status_code in (401, 403):
                fail("opencorporates: token rejected or plan exhausted.")
            resp.raise_for_status()
            body = resp.json()
    except SystemExit:
        raise
    except Exception as exc:
        fail(f"opencorporates search failed: {exc}")
    results = body.get("results", {})
    out = []
    for row in results.get("companies", []):
        c = row.get("company", {})
        out.append({
            "name": c.get("name", ""),
            "number": c.get("company_number", ""),
            "jurisdiction": c.get("jurisdiction_code", ""),
            "status": c.get("current_status", ""),
            "incorporated": c.get("incorporation_date", ""),
            "type": c.get("company_type", ""),
            "url": c.get("opencorporates_url", ""),
        })
    print(json.dumps({"query": args.query, "total": results.get("total_count", len(out)),
                      "companies": out}, ensure_ascii=False, indent=1))


# ---------------------------------------------------------------------------
# Free-source pack. No API keys anywhere in this section. Every endpoint below
# is a public government service or a free public telecom database.
# ---------------------------------------------------------------------------

AR_PARCEL_LAYER = ("https://gis.arkansas.gov/arcgis/rest/services/FEATURESERVICES"
                   "/Planning_Cadastre/FeatureServer/6/query")
CENSUS_GEOCODER = "https://geocoding.geo.census.gov/geocoder/locations/onelineaddress"


def _digits(value: str) -> str:
    return "".join(ch for ch in (value or "") if ch.isdigit())


def phone_lookup(args: argparse.Namespace) -> None:
    """Free carrier/rate-center metadata for a US or Canadian number.

    Source is LocalCallingGuide's public XML prefix API, which exposes the NANPA
    block assignment: OCN, operating company, rate center, LATA, switch CLLI and
    the rate centre coordinates. No key, no quota.

    Two honest limits. This is the ORIGINALLY ASSIGNED carrier for the NPA-NXX
    block; local number portability means the current carrier can differ. And
    there is no free CNAM source, so no subscriber name is returned - carrier
    subscriber data is CPNI-protected and not sold through any open API.
    """
    digits = _digits(args.number)
    if len(digits) == 11 and digits.startswith("1"):
        digits = digits[1:]
    if len(digits) != 10:
        fail(f"phone lookup: need a 10-digit US/Canada number, got '{args.number}'")
    npa, nxx, line = digits[:3], digits[3:6], digits[6:]

    url = f"https://localcallingguide.com/xmlprefix.php?npa={npa}&nxx={nxx}"
    try:
        with _osint_client(headers={"Accept": "text/xml,application/xml"}, timeout=30) as client:
            resp = client.get(url)
            resp.raise_for_status()
            body = resp.text
    except Exception as exc:
        fail(f"phone lookup failed: {exc}")

    import re as _re

    def tag(name: str) -> str:
        m = _re.search(rf"<{name}>(.*?)</{name}>", body, _re.DOTALL | _re.IGNORECASE)
        return (m.group(1).strip() if m else "")

    out = {
        "number": f"({npa}) {nxx}-{line}",
        "e164": "+1" + digits,
        "npa_nxx": f"{npa}-{nxx}",
        "carrier_ocn": tag("ocn"),
        "carrier": tag("company"),
        "rate_center": tag("rc"),
        "region": tag("region"),
        "lata": tag("lata"),
        "switch_clli": tag("switch"),
        "latitude": tag("lat"),
        "longitude": tag("lon"),
        "assignment_date": tag("effdate"),
        "caveat": ("Carrier shown is the original NPA-NXX block assignment. Number "
                   "portability means the current carrier may differ. No subscriber "
                   "name is available from any free source (CPNI-protected)."),
    }
    if not out["carrier"] and not out["rate_center"]:
        out["note"] = "No block assignment found. Unassigned prefix, or a toll-free/non-geographic number."
    print(json.dumps(out, ensure_ascii=False, indent=1))


def _census_geocode(client, address: str) -> tuple[float, float, str]:
    resp = client.get(CENSUS_GEOCODER, params={
        "address": address, "benchmark": "Public_AR_Current", "format": "json"})
    resp.raise_for_status()
    matches = resp.json().get("result", {}).get("addressMatches", [])
    if not matches:
        fail(f"geocode: Census returned no match for '{address}'. Try a fuller street address.")
    top = matches[0]
    coords = top.get("coordinates", {})
    return float(coords["x"]), float(coords["y"]), top.get("matchedAddress", "")


def geocode_address(args: argparse.Namespace) -> None:
    """Address to lat/long via the US Census geocoder. Free, keyless, unlimited."""
    try:
        with _osint_client(timeout=30) as client:
            lon, lat, matched = _census_geocode(client, args.address)
    except SystemExit:
        raise
    except Exception as exc:
        fail(f"geocode failed: {exc}")
    print(json.dumps({"query": args.address, "matched_address": matched,
                      "latitude": lat, "longitude": lon}, ensure_ascii=False, indent=1))


PARCEL_DB = Path(__file__).resolve().parent / "parcel_sources.db"


def _parcel_db():
    import sqlite3
    if not PARCEL_DB.exists():
        fail(f"parcel lookup: reference database missing at {PARCEL_DB}")
    conn = sqlite3.connect(str(PARCEL_DB))
    conn.row_factory = sqlite3.Row
    return conn


def _match_field(conn, role: str, available: list) -> str:
    """Pick the best field name for a role using the alias table."""
    lower = {f.lower(): f for f in available}
    for row in conn.execute(
            "SELECT pattern FROM field_alias WHERE role=? ORDER BY rank", (role,)):
        pat = row["pattern"]
        if pat in lower:
            return lower[pat]
    for row in conn.execute(
            "SELECT pattern FROM field_alias WHERE role=? ORDER BY rank", (role,)):
        pat = row["pattern"]
        for low, orig in lower.items():
            if pat in low:
                return orig
    return ""


def _introspect(conn, source, client) -> dict:
    """Read the layer's field list and cache the resolved mapping back to the db."""
    meta_url = source["query_url"].rsplit("/query", 1)[0]
    resp = client.get(meta_url, params={"f": "json"})
    resp.raise_for_status()
    fields = [f.get("name", "") for f in resp.json().get("fields", [])]
    if not fields:
        fail(f"parcel lookup: layer at {meta_url} exposed no field list")
    mapping = {role: _match_field(conn, role, fields)
               for role in ("owner", "parcel", "address", "mail", "value")}
    conn.execute("""UPDATE parcel_source SET owner_field=?, parcel_field=?,
                    address_field=?, mail_field=?, value_field=?, last_checked=date('now')
                    WHERE id=?""",
                 (mapping["owner"] or None, mapping["parcel"] or None,
                  mapping["address"] or None, mapping["mail"] or None,
                  mapping["value"] or None, source["id"]))
    conn.commit()
    return mapping


def parcel_lookup(args: argparse.Namespace) -> None:
    """Property owner of record, by street address or lat/long, for any state
    with a source registered in parcel_sources.db.

    The database stores one row per parcel service. Field names are discovered
    from the live layer on first use and cached, so adding a new state means
    inserting a URL, not reverse-engineering a schema.
    """
    lat, lon, matched, state = args.latitude, args.longitude, "", args.state.strip().upper()
    conn = _parcel_db()
    try:
        with _osint_client(timeout=45) as client:
            if args.address.strip():
                resp = client.get(CENSUS_GEOCODER, params={
                    "address": args.address, "benchmark": "Public_AR_Current", "format": "json"})
                resp.raise_for_status()
                hits = resp.json().get("result", {}).get("addressMatches", [])
                if not hits:
                    fail(f"parcel lookup: Census found no match for '{args.address}'")
                top = hits[0]
                lon = float(top["coordinates"]["x"])
                lat = float(top["coordinates"]["y"])
                matched = top.get("matchedAddress", "")
                if not state:
                    state = str(top.get("addressComponents", {}).get("state", "")).upper()
            if lat is None or lon is None:
                fail("parcel lookup: supply --address, or both --latitude and --longitude")
            if not state:
                fail("parcel lookup: could not determine the state. Pass --state (e.g. AR) "
                     "when querying by coordinates.")

            rows = conn.execute(
                "SELECT * FROM parcel_source WHERE state=? ORDER BY county IS NOT NULL, id",
                (state,)).fetchall()
            if not rows:
                covered = [r[0] for r in conn.execute(
                    "SELECT DISTINCT state FROM parcel_source ORDER BY state")]
                fail(f"parcel lookup: no source registered for {state}. "
                     f"Registered: {', '.join(covered) or 'none'}. "
                     f"Add one with: INSERT INTO parcel_source(state,name,query_url,layer_srid) "
                     f"VALUES('{state}','<name>','<.../FeatureServer/N/query>',4326); "
                     f"in {PARCEL_DB}")

            source = rows[0]
            mapping = {r: source[f"{r}_field"] for r in
                       ("owner", "parcel", "address", "mail", "value")}
            if not mapping["owner"]:
                mapping = _introspect(conn, source, client)

            def _run(src_row, field_map):
                r = client.get(src_row["query_url"], params={
                    "geometry": f"{lon},{lat}", "geometryType": "esriGeometryPoint",
                    "inSR": "4326", "spatialRel": "esriSpatialRelIntersects",
                    "outFields": "*", "returnGeometry": "false", "outSR": "4326", "f": "json"})
                r.raise_for_status()
                body = r.json()
                if body.get("error"):
                    raise RuntimeError(body["error"].get("message", "service error"))
                return body

            healed = None
            try:
                payload = _run(source, mapping)
            except Exception as dead:
                # A registered endpoint died. Arkansas did exactly this in 2015 when the
                # agency migrated servers. Retire the row, rediscover, retry once.
                conn.execute("UPDATE parcel_source SET verified=0, notes=? WHERE id=?",
                             ((source["notes"] or "") +
                              f" [RETIRED {__import__('datetime').date.today()}: {dead}]",
                              source["id"]))
                conn.commit()
                fresh = _discover_candidates(conn=conn, client=client, state=state,
                                             max_candidates=8, max_results=3)
                if not fresh:
                    fail(f"parcel lookup: registered source for {state} failed ({dead}) and "
                         f"rediscovery found no working replacement. Row marked unverified.")
                _register(conn, state, fresh[0], "Auto-replaced after previous source failed.")
                source = conn.execute(
                    "SELECT * FROM parcel_source WHERE state=? AND verified=1 "
                    "ORDER BY id DESC LIMIT 1", (state,)).fetchone()
                mapping = {r: source[f"{r}_field"] for r in
                           ("owner", "parcel", "address", "mail", "value")}
                healed = source["query_url"]
                payload = _run(source, mapping)
    except SystemExit:
        raise
    except Exception as exc:
        fail(f"parcel lookup failed: {exc}")
    finally:
        conn.close()

    if payload.get("error"):
        fail(f"parcel service error: {payload['error'].get('message', payload['error'])}")
    features = payload.get("features", [])
    if not features:
        print(json.dumps({
            "state": state, "latitude": lat, "longitude": lon, "matched_address": matched,
            "source": source["name"], "self_healed": healed, "parcels": [],
            "note": "No parcel at this point. Coverage within a state is frequently partial.",
        }, ensure_ascii=False, indent=1))
        return

    parcels = []
    for feat in features[:5]:
        a = feat.get("attributes", {})
        rec = {
            "owner": a.get(mapping["owner"], ""),
            "parcel_id": a.get(mapping["parcel"], ""),
            "property_address": a.get(mapping["address"], ""),
            "mailing_address": a.get(mapping["mail"], "") if mapping["mail"] else "",
            "assessed_value": a.get(mapping["value"], "") if mapping["value"] else "",
        }
        rec["all_attributes"] = {k: v for k, v in a.items() if v not in (None, "")}
        parcels.append(rec)

    print(json.dumps({
        "state": state, "latitude": lat, "longitude": lon, "matched_address": matched,
        "source": source["name"],
        "self_healed": healed,
        "fields_verified": bool(source["verified"]),
        "field_mapping": mapping,
        "parcel_count": len(features), "parcels": parcels,
        "caveat": ("Owner is the assessor's owner-of-record; publication lags real transfers. "
                   "Where mailing_address is blank the layer does not carry one - get it from "
                   "the county assessor. Not a legal boundary determination."),
        "source_notes": source["notes"] or "",
    }, ensure_ascii=False, indent=1))


STATE_NAMES = {
    "AL": "Alabama", "AZ": "Arizona", "AR": "Arkansas", "CA": "California",
    "CO": "Colorado", "CT": "Connecticut", "DE": "Delaware", "FL": "Florida",
    "GA": "Georgia", "ID": "Idaho", "IL": "Illinois", "IN": "Indiana",
    "IA": "Iowa", "KS": "Kansas", "KY": "Kentucky", "LA": "Louisiana",
    "ME": "Maine", "MD": "Maryland", "MA": "Massachusetts", "MI": "Michigan",
    "MN": "Minnesota", "MS": "Mississippi", "MO": "Missouri", "MT": "Montana",
    "NE": "Nebraska", "NV": "Nevada", "NH": "New Hampshire", "NJ": "New Jersey",
    "NM": "New Mexico", "NY": "New York", "NC": "North Carolina", "ND": "North Dakota",
    "OH": "Ohio", "OK": "Oklahoma", "OR": "Oregon", "PA": "Pennsylvania",
    "RI": "Rhode Island", "SC": "South Carolina", "SD": "South Dakota",
    "TN": "Tennessee", "TX": "Texas", "UT": "Utah", "VT": "Vermont",
    "VA": "Virginia", "WA": "Washington", "WV": "West Virginia",
    "WI": "Wisconsin", "WY": "Wyoming",
}

AGOL_SEARCH = "https://www.arcgis.com/sharing/rest/search"


def _probe_layer(client, conn, layer_url: str) -> dict:
    """Read a layer's fields and confirm it is publicly queryable with features."""
    meta = client.get(layer_url, params={"f": "json"}, timeout=25)
    meta.raise_for_status()
    info = meta.json()
    if info.get("error"):
        return {}
    if info.get("geometryType") != "esriGeometryPolygon":
        return {}
    fields = [f.get("name", "") for f in info.get("fields", [])]
    if not fields:
        return {}
    mapping = {role: _match_field(conn, role, fields)
               for role in ("owner", "parcel", "address", "mail", "value")}
    if not mapping["owner"] or not mapping["parcel"]:
        return {}
    count = client.get(layer_url + "/query", timeout=25, params={
        "where": "1=1", "returnCountOnly": "true", "f": "json"})
    count.raise_for_status()
    body = count.json()
    if body.get("error"):
        return {}
    n = body.get("count", 0)
    if n < 1000:
        return {}
    srid = (info.get("extent", {}).get("spatialReference", {}) or {}).get("wkid", 4326)
    return {"layer_url": layer_url, "name": info.get("name", ""), "count": n,
            "srid": srid, "mapping": mapping}


def _discover_candidates(client, conn, state: str, max_candidates: int = 12,
                         max_results: int = 5) -> list:
    """Search ArcGIS Online for a state's parcel layers and return only verified ones."""
    full = STATE_NAMES[state]
    seen, candidates, results = set(), [], []
    queries = [
        f'"{full}" AND parcels AND type:"Feature Service"',
        f'"{full}" AND "tax parcels" AND type:"Feature Service"',
        f'"{full} statewide parcels" AND type:"Feature Service"',
        f'"{full}" AND cadastral AND type:"Feature Service"',
        f'"{full}" AND "land records" AND type:"Feature Service"',
    ]
    for q in queries:
        try:
            resp = client.get(AGOL_SEARCH, params={
                "f": "json", "q": q, "num": 20,
                "sortField": "numViews", "sortOrder": "desc"})
            resp.raise_for_status()
            for item in resp.json().get("results", []):
                url = (item.get("url") or "").rstrip("/")
                if url and url not in seen:
                    seen.add(url)
                    candidates.append((item.get("title", ""), item.get("owner", ""), url))
        except Exception:
            continue

    for title, owner, svc in candidates[:max_candidates]:
        if "/FeatureServer" not in svc and "/MapServer" not in svc:
            continue
        try:
            top = client.get(svc, params={"f": "json"}, timeout=25)
            top.raise_for_status()
            layers = top.json().get("layers", [])
        except Exception:
            continue
        for lyr in layers[:12]:
            lid = lyr.get("id")
            if lid is None:
                continue
            if "parcel" not in str(lyr.get("name", "")).lower() and len(layers) > 1:
                continue
            try:
                probe = _probe_layer(client, conn, f"{svc}/{lid}")
            except Exception:
                continue
            if probe:
                probe.update({"item_title": title, "item_owner": owner})
                results.append(probe)
        if len(results) >= max_results:
            break
    results.sort(key=lambda r: r["count"], reverse=True)
    return results


def _register(conn, state: str, best: dict, note_prefix: str) -> None:
    m = best["mapping"]
    conn.execute("""INSERT INTO parcel_source
        (state,county,name,query_url,layer_srid,owner_field,parcel_field,
         address_field,mail_field,value_field,verified,notes,last_checked)
        VALUES (?,NULL,?,?,?,?,?,?,?,?,1,?,date('now'))""",
        (state, f"{best['item_title']} - {best['name']}",
         best["layer_url"] + "/query", best["srid"],
         m["owner"] or None, m["parcel"] or None, m["address"] or None,
         m["mail"] or None, m["value"] or None,
         f"{note_prefix} Publisher: {best['item_owner']}. "
         f"{best['count']} features at registration."))
    conn.commit()


def parcel_discover(args: argparse.Namespace) -> None:
    """Find and verify a public parcel service for a state, then optionally register it."""
    state = args.state.strip().upper()
    if state not in STATE_NAMES:
        fail(f"discover: '{args.state}' is not a continental US state code")
    conn = _parcel_db()
    committed = None
    try:
        with _osint_client(timeout=30) as client:
            results = _discover_candidates(conn=conn, client=client, state=state,
                                           max_candidates=args.max_candidates,
                                           max_results=args.max_results)
        if args.commit and results:
            _register(conn, state, results[0], "Auto-discovered via ArcGIS Online.")
            committed = results[0]["layer_url"]
    finally:
        conn.close()
    print(json.dumps({
        "state": state, "state_name": STATE_NAMES[state],
        "verified_layers": [{
            "layer_url": r["layer_url"], "layer": r["name"], "item": r["item_title"],
            "publisher": r["item_owner"], "features": r["count"], "srid": r["srid"],
            "fields": r["mapping"],
        } for r in results[:args.max_results]],
        "committed": committed,
        "note": ("Nothing was written." if not args.commit else
                 "Best candidate registered." if committed else
                 "Nothing verified, so nothing written."),
    }, ensure_ascii=False, indent=1))


def parcel_sources(args: argparse.Namespace) -> None:
    """List what is currently registered in the parcel reference database."""
    conn = _parcel_db()
    try:
        rows = conn.execute(
            "SELECT state,county,name,query_url,verified,last_checked,notes "
            "FROM parcel_source ORDER BY state, county IS NOT NULL, id").fetchall()
    finally:
        conn.close()
    missing = sorted(set(STATE_NAMES) - {r["state"] for r in rows})
    print(json.dumps({
        "registered": [dict(r) for r in rows],
        "states_covered": sorted({r["state"] for r in rows}),
        "states_missing": missing,
        "hint": "Register a missing state with the parcel_source_discover tool.",
    }, ensure_ascii=False, indent=1))


def _entity_recipe(conn, state: str):
    row = conn.execute("SELECT * FROM entity_source WHERE state=? ORDER BY verified DESC, id",
                       (state,)).fetchone()
    if not row:
        have = [r[0] for r in conn.execute("SELECT DISTINCT state FROM entity_source ORDER BY state")]
        fail(f"entity search: no registry recipe for {state}. Registered: {', '.join(have) or 'none'}. "
             f"Recipes live in the entity_source table of {PARCEL_DB}.")
    return row, json.loads(row["recipe"])


def _pick(row: dict, keys: list) -> str:
    """Pull the first present key from a result row; keys may be names or list indices."""
    if isinstance(row, dict):
        lower = {str(k).lower(): v for k, v in row.items()}
        for k in keys:
            if str(k).lower() in lower and lower[str(k).lower()] not in (None, ""):
                return str(lower[str(k).lower()])
    elif isinstance(row, list):
        for k in keys:
            if str(k).isdigit() and int(k) < len(row):
                val = row[int(k)]
                if val not in (None, ""):
                    return _compact_text(str(val), 200)
    return ""


def entity_search(args: argparse.Namespace) -> None:
    """Look up a business entity in a state registry: status, type, filing number,
    and registered agent.

    Built for the case where a parcel's owner of record is an LLC or trust rather
    than a person. The registered agent is the entity's official address for
    service and is usually the right contact.
    """
    state = args.state.strip().upper() or "AR"
    if not args.name.strip() and not args.agent.strip():
        fail("entity search: supply --name or --agent")
    conn = _parcel_db()
    try:
        row, rec = _entity_recipe(conn, state)
    finally:
        conn.close()

    try:
        with _osint_client(headers={"Accept": "application/json, text/html"}, timeout=40) as client:
            payload = dict(rec.get("static_params", {}))
            pmap = rec.get("params", {})
            if args.name.strip() and pmap.get("name"):
                payload[pmap["name"]] = args.name.strip()
            if args.agent.strip() and pmap.get("agent"):
                payload[pmap["agent"]] = args.agent.strip()

            csrf = rec.get("csrf") or {}
            if csrf.get("source") == "meta":
                page = client.get(rec["form_url"])
                page.raise_for_status()
                import re as _re
                m = _re.search(
                    rf'<meta[^>]+name=["\']{csrf.get("meta_name","csrf-token")}["\'][^>]+content=["\']([^"\']+)',
                    page.text, _re.IGNORECASE)
                if not m:
                    m = _re.search(
                        rf'<meta[^>]+content=["\']([^"\']+)["\'][^>]+name=["\']{csrf.get("meta_name","csrf-token")}["\']',
                        page.text, _re.IGNORECASE)
                if m and csrf.get("param"):
                    payload[csrf["param"]] = m.group(1)

            method = rec.get("method", "POST").upper()
            if method == "POST":
                resp = client.post(rec["search_url"], data=payload,
                                   headers={"X-Requested-With": "XMLHttpRequest",
                                            "Referer": rec["form_url"]})
            else:
                resp = client.get(rec["search_url"], params=payload,
                                  headers={"Referer": rec["form_url"]})
            resp.raise_for_status()
            text = resp.text
    except Exception as exc:
        fail(f"entity search failed against {state} registry: {exc}. "
             f"If this is a 404/419/422 the stored request recipe is wrong - see the notes "
             f"column of entity_source for how to capture the correct one.")

    rows, fmt = [], "unknown"
    try:
        body = json.loads(text)
        fmt = "json"
        path = rec.get("result", {}).get("rows_path", "data")
        cur = body
        for part in str(path).split("."):
            if isinstance(cur, dict) and part in cur:
                cur = cur[part]
        rows = cur if isinstance(cur, list) else []
    except ValueError:
        fmt = "html"
        import re as _re
        for tr in _re.findall(r"<tr[^>]*>(.*?)</tr>", text, _re.DOTALL | _re.IGNORECASE):
            cells = [_compact_text(_re.sub(r"<[^>]+>", " ", td), 200)
                     for td in _re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", tr, _re.DOTALL | _re.IGNORECASE)]
            if any(cells):
                rows.append(cells)
        if rows and all(not c or c.lower() in
                        ("name", "filing #", "status", "type", "registered agent")
                        for c in rows[0]):
            rows = rows[1:]

    fields = rec.get("result", {}).get("fields", {})
    out = []
    for r in rows[:max(1, min(args.max_results, 30))]:
        rec_out = {k: _pick(r, v) for k, v in fields.items()}
        if any(rec_out.values()):
            out.append(rec_out)

    if not out:
        fail(f"entity search: request succeeded ({fmt}) but no rows parsed. The stored recipe's "
             f"result mapping does not match what this registry returned. Capture one real "
             f"request/response and UPDATE entity_source for {state}.")

    print(json.dumps({
        "state": state, "registry": row["name"],
        "recipe_verified": bool(row["verified"]),
        "query": {"name": args.name, "agent": args.agent},
        "result_count": len(out), "entities": out,
        "note": ("Registered agent is the entity's official address for service of process "
                 "and is generally the correct contact. It is often a law firm or commercial "
                 "agent rather than a principal."),
    }, ensure_ascii=False, indent=1))


def broker_registry(args: argparse.Namespace) -> None:
    """Search California's data broker registry - who collects and sells personal data.

    Every data broker doing business in California must register annually and
    disclose what categories it collects and how to opt out. The registry is the
    authoritative list of who holds consumer data, and it is public and free.
    Useful from any state: the brokers are national, and most run one opt-out
    pipeline rather than maintaining per-state ones.
    """
    import csv
    import io
    import re as _re

    page = "https://cppa.ca.gov/data_broker_registry/"
    try:
        with _osint_client(headers={"Accept": "text/html,text/csv,*/*"}, timeout=60) as client:
            resp = client.get(page)
            resp.raise_for_status()
            links = _re.findall(r'href=["\']([^"\']+\.csv[^"\']*)["\']', resp.text, _re.IGNORECASE)
            if not links:
                fail("broker registry: no CSV link found on the CPPA registry page. "
                     "The page layout changed - open " + page + " and grab the Download CSV link.")
            target = links[0]
            if target.startswith("/"):
                target = "https://cppa.ca.gov" + target
            data = client.get(target)
            data.raise_for_status()
            text = data.content.decode("utf-8-sig", errors="replace")
    except SystemExit:
        raise
    except Exception as exc:
        fail(f"broker registry fetch failed: {exc}")

    rows = list(csv.DictReader(io.StringIO(text)))
    if not rows:
        fail("broker registry: CSV parsed empty")

    needle = args.query.strip().lower()
    matches = rows if not needle else [
        r for r in rows if any(needle in str(v).lower() for v in r.values())]

    limit = max(1, min(args.max_results, 50))
    out = []
    for r in matches[:limit]:
        rec = {}
        for k, v in r.items():
            if v and str(v).strip():
                rec[k.strip()] = str(v).strip()[:400]
        out.append(rec)
    print(json.dumps({
        "source": page,
        "total_registered": len(rows),
        "matched": len(matches),
        "returned": len(out),
        "brokers": out,
    }, ensure_ascii=False, indent=1))


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(prog="helm_tools")
    sub = parser.add_subparsers(dest="command", required=True)

    p = sub.add_parser("web-search")
    p.add_argument("--query", required=True)
    p.add_argument("--max-results", type=int, default=8)
    p.set_defaults(func=web_search)

    p = sub.add_parser("github-search")
    p.add_argument("--query", required=True)
    p.add_argument("--max-results", type=int, default=15)
    p.add_argument("--sort", default="best-match")
    p.set_defaults(func=github_search)

    p = sub.add_parser("seen-list")
    p.add_argument("--file", required=True)
    p.add_argument("--add", default="")
    p.add_argument("--limit", type=int, default=500)
    p.set_defaults(func=seen_list)

    p = sub.add_parser("web-fetch")
    p.add_argument("--url", required=True)
    p.add_argument("--max-chars", type=int, default=12000)
    p.set_defaults(func=web_fetch)

    p = sub.add_parser("extract-document")
    p.add_argument("--path", required=True)
    p.add_argument("--max-chars", type=int, default=200000)
    p.set_defaults(func=extract_document)

    p = sub.add_parser("desktop-screenshot")
    p.add_argument("--output", required=True)
    p.add_argument("--monitor", type=int, default=0)
    p.set_defaults(func=desktop_screenshot)

    p = sub.add_parser("desktop-click")
    p.add_argument("--x", type=int, required=True)
    p.add_argument("--y", type=int, required=True)
    p.add_argument("--button", choices=["left", "right", "middle"], default="left")
    p.add_argument("--clicks", type=int, default=1)
    p.add_argument("--interval-ms", type=int, default=100)
    p.set_defaults(func=desktop_click)

    p = sub.add_parser("desktop-type")
    p.add_argument("--text", required=True)
    p.add_argument("--interval-ms", type=int, default=10)
    p.set_defaults(func=desktop_type)

    p = sub.add_parser("desktop-hotkey")
    p.add_argument("--keys", required=True)
    p.set_defaults(func=desktop_hotkey)

    p = sub.add_parser("comfy-generate")
    p.add_argument("--url", required=True)
    p.add_argument("--workflow", required=True)
    p.add_argument("--prompt", required=True)
    p.add_argument("--negative-prompt", default="")
    p.add_argument("--width", type=int, default=1024)
    p.add_argument("--height", type=int, default=1024)
    p.add_argument("--output-dir", required=True)
    p.add_argument("--timeout-seconds", type=int, default=900)
    p.add_argument("--start-command", default="")
    p.set_defaults(func=comfy_generate)

    p = sub.add_parser("rdap")
    p.add_argument("--target", required=True)
    p.set_defaults(func=rdap_lookup)

    p = sub.add_parser("dns-lookup")
    p.add_argument("--domain", required=True)
    p.add_argument("--types", default="A")
    p.set_defaults(func=dns_lookup)

    p = sub.add_parser("crtsh")
    p.add_argument("--domain", required=True)
    p.add_argument("--max-results", type=int, default=200)
    p.add_argument("--timeout", type=int, default=90)
    p.set_defaults(func=crtsh_subdomains)

    p = sub.add_parser("edgar-company")
    p.add_argument("--query", required=True)
    p.add_argument("--max-results", type=int, default=15)
    p.set_defaults(func=edgar_company)

    p = sub.add_parser("edgar-search")
    p.add_argument("--query", required=True)
    p.add_argument("--forms", default="")
    p.add_argument("--max-results", type=int, default=15)
    p.set_defaults(func=edgar_search)

    p = sub.add_parser("hibp-account")
    p.add_argument("--account", required=True)
    p.set_defaults(func=hibp_account)

    p = sub.add_parser("opencorporates")
    p.add_argument("--query", required=True)
    p.add_argument("--jurisdiction", default="")
    p.add_argument("--max-results", type=int, default=15)
    p.set_defaults(func=opencorporates_search)


    p = sub.add_parser("phone-lookup")
    p.add_argument("--number", required=True)
    p.set_defaults(func=phone_lookup)

    p = sub.add_parser("geocode")
    p.add_argument("--address", required=True)
    p.set_defaults(func=geocode_address)

    p = sub.add_parser("parcel-lookup")
    p.add_argument("--address", default="")
    p.add_argument("--state", default="")
    p.add_argument("--latitude", type=float, default=None)
    p.add_argument("--longitude", type=float, default=None)
    p.set_defaults(func=parcel_lookup)

    p = sub.add_parser("broker-registry")
    p.add_argument("--query", default="")
    p.add_argument("--max-results", type=int, default=20)
    p.set_defaults(func=broker_registry)


    p = sub.add_parser("parcel-discover")
    p.add_argument("--state", required=True)
    p.add_argument("--commit", action="store_true")
    p.add_argument("--max-candidates", type=int, default=12)
    p.add_argument("--max-results", type=int, default=5)
    p.set_defaults(func=parcel_discover)

    p = sub.add_parser("parcel-sources")
    p.set_defaults(func=parcel_sources)


    p = sub.add_parser("entity-search")
    p.add_argument("--name", default="")
    p.add_argument("--agent", default="")
    p.add_argument("--state", default="AR")
    p.add_argument("--max-results", type=int, default=15)
    p.set_defaults(func=entity_search)

    return parser


def _force_utf8_streams() -> None:
    """Make stdout/stderr UTF-8 regardless of the Windows code page.

    The C++ host reads this script's stdout through a pipe. Python sees a pipe,
    not a console, so it picks the locale encoding - cp1252 on a US Windows box -
    and any non-ASCII character in fetched web text kills the process:

        UnicodeEncodeError: 'charmap' codec can't encode character '\u2010'

    U+2010 is an ordinary typographic hyphen. Wikipedia and half the web are full
    of them, so this is not an edge case; it is every other page. nlohmann::json
    on the receiving side expects UTF-8, which is what this produces.
    """
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass


def main() -> None:
    _force_utf8_streams()
    args = build_parser().parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
